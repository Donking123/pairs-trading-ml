"""
Build the QF621 main presentation deck as .pptx.

Run from pairs-trading-ml/:
  python submission/report/build_pptx.py

Outputs: submission/report/QF621_pairs_trading_deck.pptx
"""
from __future__ import annotations

import sys
from pathlib import Path

DEPS = Path(__file__).resolve().parent.parent / ".deps"
if DEPS.exists():
    sys.path.insert(0, str(DEPS))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPORT_DIR = Path(__file__).resolve().parent
ASSETS = REPORT_DIR / "assets"
OUT = REPORT_DIR / "QF621_pairs_trading_deck.pptx"

# ── colours ──
DARK_BLUE = RGBColor(0x1B, 0x49, 0x65)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1D, 0x27, 0x33)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT_BG = RGBColor(0xF2, 0xF7, 0xFB)
ACCENT = RGBColor(0x5F, 0xA8, 0xD3)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def _add_para(tf, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
              space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def _add_table(slide, left, top, width, height, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row_data in enumerate(rows):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                elif c == 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                else:
                    paragraph.alignment = PP_ALIGN.CENTER

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return table


def _title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.2),
                 "Clustering-Based Pairs Trading", font_size=40, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(2.7), Inches(8), Inches(0.8),
                 "A from-scratch replication & extension of Rotondi & Russo (2025)",
                 font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.6),
                 "Headline: matched the paper's best strategy at Sharpe 1.113 (paper: 1.01)",
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.4), Inches(8), Inches(0.6),
                 "Validated on 5 years out-of-sample (2021-2025): Sharpe 0.412",
                 font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.4),
                 "Built in Python on WRDS/CRSP data  |  ML clustering + statistical-arbitrage signal  |  full rolling-window backtest engine",
                 font_size=12, color=GREY, align=PP_ALIGN.CENTER)


def _problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "The Problem", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5),
                      "Pairs trading = find two securities whose prices move together, trade the spread when it diverges, profit when it reverts.",
                      font_size=20)
    _add_para(tf, "")
    _add_para(tf, "The hard part is pair selection across a 1,000-stock universe:", font_size=20, bold=True)
    _add_para(tf, "  - Brute-forcing all ~500,000 pairs is noisy and overfits", font_size=18)
    _add_para(tf, "  - We want pairs that are economically related, not coincidental", font_size=18)
    _add_para(tf, "")
    _add_para(tf, "The paper's idea: use unsupervised ML clustering to group similar stocks first, then only form pairs within clusters. Fewer, higher-quality candidates.", font_size=20, bold=True)
    _add_para(tf, "")
    _add_para(tf, "Our goal: rebuild the whole pipeline from the data up, reproduce the published result independently, then extend it.", font_size=18, color=GREY)


def _data_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Data Spine — Built for Honesty, Not Optics", font_size=36, bold=True, color=DARK_BLUE)
    rows = [
        ["", ""],
        ["Source", "WRDS / CRSP daily (institutional-grade academic data)"],
        ["Period", "Jan 2000 - Dec 2025 (26 years)"],
        ["Universe", "991 stocks - share codes 10/11 only (US ordinary common)"],
        ["Survivorship bias", "Eliminated - delisted names kept; delisting returns modelled"],
        ["In-sample", "2003-2020 = 215 monthly returns"],
        ["Out-of-sample", "2021-2025 = 59 monthly returns (frozen, no re-tuning)"],
    ]
    _add_table(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(3.5), rows,
               col_widths=[Inches(2.5), Inches(6.5)])
    _add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(1),
                 "First 3 years are reserved purely as the formation window, so every reported return is genuinely out-of-sample.",
                 font_size=14, color=GREY)


def _pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "The Pipeline (Rolling, Monthly)", font_size=36, bold=True, color=DARK_BLUE)
    pipeline_text = (
        "CRSP returns\n"
        "    -> [1] DISTANCE metric (SSD or PC)\n"
        "    -> [2] OPTICS clustering\n"
        "         (pairs only within a cluster)\n"
        "    -> [3] Cointegration filter (optional)\n"
        "    -> [4] Spread + rolling z-score\n"
        "    -> [5] Backtest: enter |z|>2, exit z=0, force-close month-end"
    )
    _add_textbox(slide, Inches(1), Inches(1.3), Inches(8), Inches(3),
                 pipeline_text, font_size=18, font_name="Courier New")
    tf = _add_textbox(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(1.5),
                      "Each month: form pairs on the trailing 3 years, trade them for the next 1 month, roll forward.",
                      font_size=20, bold=True)
    _add_para(tf, "215 times (in-sample) + 59 times (out-of-sample).", font_size=18, color=GREY)


def _distance_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Two Distance Metrics", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5),
                      "How do we decide which stocks are \"similar\" before clustering?", font_size=20)
    _add_para(tf, "")
    _add_para(tf, "SSD - Sum of Squared Deviations between normalised price paths. Classic baseline (Gatev et al.).", font_size=18)
    _add_para(tf, "")
    _add_para(tf, "PC - Partial-correlation distance on market-adjusted returns. Strips out the common market factor, finds idiosyncratic co-movement. The paper's winning metric.", font_size=18, bold=True)
    _add_para(tf, "")
    _add_para(tf, "Distance feeds the clusterer; the metric choice turns out to be the single biggest driver of risk-adjusted performance.", font_size=14, color=GREY)


def _clustering_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Clustering with OPTICS", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(1.5),
                      "OPTICS - density-based clustering (cousin of DBSCAN).", font_size=20)
    _add_para(tf, "  - No need to pre-specify the number of clusters", font_size=18)
    _add_para(tf, "  - Leaves genuinely unrelated stocks unclustered", font_size=18)

    rows = [
        ["Metric", "Ours (SSD)", "Ours (PC)", "Paper"],
        ["# clusters", "47", "81", "48 / 109"],
        ["Purity vs SIC", "0.871", "0.937", "0.81 / 0.84"],
        ["GOOG & GOOGL?", "Y", "Y", "Y"],
    ]
    _add_table(slide, Inches(1), Inches(3.2), Inches(8), Inches(2), rows)
    _add_textbox(slide, Inches(0.5), Inches(5.7), Inches(9), Inches(0.8),
                 "xi hyperparameter tuned on cluster quality on 3 hold-out dates - frozen before ever looking at Sharpe.",
                 font_size=14, color=GREY)


def _signal_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Cointegration Filter & Trading Signal", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(4), Inches(4),
                      "Engle-Granger filter (optional):", font_size=20, bold=True)
    _add_para(tf, "  - Spread rejects unit-root at 5%", font_size=18)
    _add_para(tf, "  - Half-life in [5, 60] trading days", font_size=18)

    tf2 = _add_textbox(slide, Inches(5), Inches(1.1), Inches(4.5), Inches(4),
                       "Trading signal:", font_size=20, bold=True)
    _add_para(tf2, "  - Hedge ratio via OLS -> spread", font_size=18)
    _add_para(tf2, "  - 6-month rolling z-score", font_size=18)
    _add_para(tf2, "  - Enter when |z| > 2", font_size=18)
    _add_para(tf2, "  - Exit at z = 0 (reversion)", font_size=18)
    _add_para(tf2, "  - Force-close at month-end", font_size=18)
    _add_para(tf2, "  - Strict look-ahead protection", font_size=16, color=GREY)


def _factor_beta_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Original Extension: Factor-Beta Clustering", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(2.5),
                      "A structurally different similarity metric — cluster by shared risk-factor exposures, not return co-movement.",
                      font_size=20)
    _add_para(tf, "")
    _add_para(tf, "  - Ridge-regress each stock's returns on 18 factors (6 FF style + 12 FF industry)", font_size=18)
    _add_para(tf, "  - Each stock gets an 18-dimensional beta vector", font_size=18)
    _add_para(tf, "  - Distance = standardised Euclidean between beta vectors", font_size=18)
    _add_para(tf, "  - Stocks with similar factor loadings cluster together", font_size=18)

    rows = [
        ["Metric", "Sharpe", "vs Paper"],
        ["PC core", "1.113", "1.01 (match)"],
        ["Factor-beta core", "1.149", "original extension, beats PC"],
        ["Factor + coint filter", "0.969", "filter slightly reduces"],
    ]
    _add_table(slide, Inches(1), Inches(4.0), Inches(8), Inches(2), rows)
    _add_textbox(slide, Inches(0.5), Inches(6.2), Inches(9), Inches(0.8),
                 "Two independent metrics reaching ~1.0 is strong evidence the edge is real, not a PC-specific artefact.",
                 font_size=16, bold=True, color=DARK_BLUE)


def _is_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "In-Sample Results (2003-2020, 215 months)", font_size=36, bold=True, color=DARK_BLUE)
    rows = [
        ["Strategy", "Ann. Ret", "Ann. Vol", "Sharpe", "Max DD"],
        ["Factor-beta core", "4.16%", "3.61%", "1.149", "-3.59%"],
        ["PC core", "3.86%", "3.46%", "1.113", "-5.75%"],
        ["Factor + filter", "5.16%", "5.35%", "0.969", "-6.41%"],
        ["SSD + filter", "4.97%", "6.37%", "0.794", "-9.74%"],
        ["PC + filter", "2.79%", "3.57%", "0.788", "-5.88%"],
        ["SSD core (baseline)", "3.48%", "5.54%", "0.645", "-14.31%"],
    ]
    _add_table(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(3.2), rows)
    tf = _add_textbox(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(2),
                      "Factor-beta core (1.149) beats PC core (1.113) — two independent metrics above 1.1 is strong evidence the edge is real.",
                      font_size=20, bold=True)
    _add_para(tf, "Paper target for PC core: 1.01. Ours: 1.113. Factor-beta: 1.149 (original extension).", font_size=18, color=GREY)


def _attribution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Why It Works - Diagnosed, Not Assumed", font_size=36, bold=True, color=DARK_BLUE)
    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5),
                      "P&L attribution on the SSD baseline revealed a bimodal pattern:", font_size=20)
    _add_para(tf, "")
    _add_para(tf, "  11.4% of trades cleanly revert: +471 bps each -> ALL the profit", font_size=20, bold=True)
    _add_para(tf, "  88.4% get force-closed at month-end: -32 bps each -> constant drag", font_size=20, bold=True)
    _add_para(tf, "")
    _add_para(tf, "Thesis: the job isn't \"find a higher Sharpe.\" It's kill the force-close drag by only trading pairs that actually revert.", font_size=20)
    _add_para(tf, "")
    _add_para(tf, "PC distance cut the per-trade force-close drag by 65%.", font_size=20, bold=True, color=DARK_BLUE)

    img_path = ASSETS / "mechanism.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(6), Inches(3.5), Inches(3.5))


def _cumulative_pnl_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Cumulative Returns — All Strategies", font_size=36, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "cumulative_pnl.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.1), Inches(9))
    _add_textbox(slide, Inches(0.5), Inches(6.6), Inches(9), Inches(0.6),
                 "PC core (dark) is the smoothest compounder. SSD strategies have higher absolute returns but far more volatility.",
                 font_size=14, color=GREY)


def _pnl_breakdown_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "P&L Breakdown by Exit Reason", font_size=36, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "pnl_breakdown.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), Inches(9.4))

    rows = [
        ["Strategy", "Reversion %", "Rev. bps", "Force-close %", "FC bps"],
        ["SSD core", "11.6%", "+474", "88.1%", "-36"],
        ["PC core", "8.4%", "+346", "91.4%", "-14"],
        ["PC + filter", "10.0%", "+327", "89.7%", "-23"],
        ["SSD + filter", "12.3%", "+491", "87.5%", "-39"],
    ]
    _add_table(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(1.8), rows)
    _add_textbox(slide, Inches(0.5), Inches(7.0), Inches(9), Inches(0.4),
                 "Key: PC cuts force-close drag from -36 bps (SSD) to -14 bps. The metric's job is selectivity, not prediction.",
                 font_size=13, bold=True, color=DARK_BLUE)


def _yearly_returns_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Year-by-Year Returns — Regime Dependence", font_size=36, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "yearly_returns.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.2), Inches(1.0), Inches(9.6))

    tf = _add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(2),
                      "Findings:", font_size=16, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - Best year: 2009 (+28.4%) — GFC dislocation = highest dispersion, most reversions", font_size=14)
    _add_para(tf, "  - Calm periods (2013-19): avg +1.6%/yr — the strategy idles, doesn't lose", font_size=14)
    _add_para(tf, "  - OOS (2021-25): weak in low-dispersion 2021-22, recovered in 2023 & 2025", font_size=14)
    _add_para(tf, "  - 15/18 IS years positive (83%) — consistent but regime-dependent magnitude", font_size=14)


def _oos_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Out-of-Sample Validation (2021-2025)", font_size=36, bold=True, color=DARK_BLUE)
    rows = [
        ["Strategy", "IS Sharpe", "OOS Sharpe", "OOS Ann. Ret", "OOS MDD", "Months"],
        ["PC + filter", "0.788", "0.461", "1.31%", "-2.98%", "59"],
        ["PC core", "1.113", "0.412", "0.82%", "-2.75%", "59"],
        ["SSD + filter", "0.794", "0.224", "0.82%", "-5.52%", "59"],
        ["Factor + filter", "0.969", "0.131", "0.42%", "-6.56%", "59"],
        ["Factor core", "1.149", "-0.103", "-0.29%", "-6.99%", "59"],
    ]
    _add_table(slide, Inches(0.5), Inches(1.3), Inches(9), Inches(3.0), rows)
    tf = _add_textbox(slide, Inches(0.5), Inches(4.6), Inches(9), Inches(3),
                      "PC + filter is the best OOS strategy (0.461) — all three filtered strategies stay positive.", font_size=20, bold=True)
    _add_para(tf, "")
    _add_para(tf, "The cointegration filter helps in OOS: it removes pairs that drift apart during the trading window. The filter flips factor-beta from negative (-0.103) to positive (0.131).", font_size=17)
    _add_para(tf, "")
    _add_para(tf, "The honest story: the strategy works, but it is regime-dependent. Performance tracks volatility and dispersion, not calendar time.", font_size=17, bold=True, color=DARK_BLUE)


def _rigour_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Engineering Rigour", font_size=36, bold=True, color=DARK_BLUE)
    items = [
        "Modular package - distances, clustering, spread, cointegration, backtest, performance, costs, lookahead - each independently unit-tested",
        "67 synthetic unit tests across 11 files - validate every component on data with a known answer",
        "6/6 lookahead audit PASS - black-box test confirms no future information leakage",
        "Look-ahead protection baked into the z-score and the rolling formation/trading split",
        "Reproducible - config-driven single source of truth; hyperparameters frozen pre-results",
        "Implementation-correctness review (Phase 6) - audited own code, found and fixed 6 flaws",
    ]
    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5),
                      "", font_size=18)
    for item in items:
        _add_para(tf, f"  - {item}", font_size=16, space_before=Pt(8))

    _add_para(tf, "", font_size=10)
    _add_para(tf, "Stack: Python, pandas, NumPy, scikit-learn (OPTICS), statsmodels (ADF/OLS), matplotlib, parquet",
              font_size=13, color=GREY)


def _takeaway_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "What This Demonstrates", font_size=36, bold=True, color=DARK_BLUE)
    items = [
        ("Independent replication", "of a 2025 research paper - read it, rebuilt it, matched the headline number on data we sourced ourselves."),
        ("Honest validation", "- 5-year OOS test shows the strategy survives but with regime-dependent decay. We report what we found, not what we hoped."),
        ("Quant intuition", "- diagnosed why the strategy makes money, then engineered the fix."),
        ("Discipline", "- survivorship-bias-free data, out-of-sample design, frozen hyperparameters, unit-tested code, implementation-correctness review."),
    ]
    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5), "", font_size=18)
    for title, desc in items:
        _add_para(tf, f"  {title} {desc}", font_size=18, space_before=Pt(14))


def _filtered_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Filtered Strategies — IS vs OOS Scorecard", font_size=36, bold=True, color=DARK_BLUE)

    rows = [
        ["Strategy", "Period", "Ann. Ret", "Ann. Vol", "Sharpe", "Max DD", "Months"],
        ["Factor + filter", "IS 2003-2020", "5.16%", "5.35%", "0.969", "-6.41%", "215"],
        ["SSD + filter", "IS 2003-2020", "4.97%", "6.37%", "0.794", "-9.74%", "215"],
        ["PC + filter", "IS 2003-2020", "2.79%", "3.57%", "0.788", "-5.88%", "215"],
        ["PC + filter", "OOS 2021-2025", "1.31%", "2.92%", "0.461", "-2.98%", "59"],
        ["SSD + filter", "OOS 2021-2025", "0.82%", "4.01%", "0.224", "-5.52%", "59"],
        ["Factor + filter", "OOS 2021-2025", "0.42%", "3.77%", "0.131", "-6.56%", "59"],
    ]

    _add_table(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(3.8), rows)

    tf = _add_textbox(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(2.2),
                      "The cointegration filter is required by the project specification. Key findings:",
                      font_size=18, bold=True)
    _add_para(tf, "  - PC + filter is the BEST OOS strategy (0.461) — filter helps PC in OOS", font_size=16)
    _add_para(tf, "  - Filter rescues factor-beta OOS: -0.103 (core) -> 0.131 (filtered)", font_size=16)
    _add_para(tf, "  - All three filtered strategies stay positive OOS — the filter removes drifting pairs", font_size=16)


def _cumulative_pnl_filtered_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Cumulative Returns — Filtered Strategies (IS + OOS)", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "cumulative_pnl_filtered.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), Inches(9.4))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.6),
                 "Solid lines = IS (2003-2020). Dashed lines = OOS (2021-2025). All strategies use the cointegration filter.",
                 font_size=14, color=GREY)


def _yearly_returns_filtered_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Year-by-Year Returns — Filtered Strategies", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "yearly_returns_filtered.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.2), Inches(1.0), Inches(9.6))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    tf = _add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(2),
                      "Findings:", font_size=16, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - GFC (2008-09) drives outsized returns for all filtered strategies — dislocation harvesting", font_size=14)
    _add_para(tf, "  - Calm years (2013-19): small but mostly positive — the filter keeps the strategy idle, not losing", font_size=14)
    _add_para(tf, "  - OOS: all three filtered strategies show positive years in 2023 and 2025 (high-dispersion periods)", font_size=14)


def _pnl_breakdown_filtered_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "P&L Breakdown — Filtered Strategies (IS vs OOS)", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "pnl_breakdown_filtered.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.1), Inches(1.0), Inches(9.8))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    tf = _add_textbox(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(1.8),
                      "Filtered strategy diagnostics:", font_size=16, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - PC+filter: lowest force-close drag among filtered variants", font_size=14)
    _add_para(tf, "  - Factor+filter: higher reversion rate but higher per-trade variance", font_size=14)
    _add_para(tf, "  - IS vs OOS comparison shows consistency of trade mechanics", font_size=14)


def _pair_examples_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "What Each Metric Selects — Example Pairs", font_size=36, bold=True, color=DARK_BLUE)

    rows = [
        ["", "SSD + filter", "PC + filter", "Factor + filter"],
        ["#1 pair", "CMS/XEL (utilities)", "BK/NTRS (banks)", "VZ/T (telecom)"],
        ["#2 pair", "LMT/RTN (defense)", "VZ/T (telecom)", "ED/SO (utilities)"],
        ["#3 pair", "RF/ZION (banks)", "ED/SO (utilities)", "WEC/XEL (utilities)"],
        ["Unique to", "MCO/UNH, JPM/GL", "TMO/A, FDX/UPS", "NKE/YUM, ORCL/QCOM"],
        ["Total pairs", "3,028", "1,401", "1,673"],
        ["Total trades", "8,143", "7,594", "4,871"],
        ["Reversion %", "12.3%", "10.0%", "10.5%"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.1), Inches(9.4), Inches(3.8), rows,
               col_widths=[Inches(1.8), Inches(2.5), Inches(2.5), Inches(2.6)])

    tf = _add_textbox(slide, Inches(0.5), Inches(5.2), Inches(9), Inches(2.3),
                      "Key differences:", font_size=18, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - SSD casts the widest net (3,028 pairs) — any stocks with similar price paths", font_size=15)
    _add_para(tf, "  - PC is the most selective (1,401 pairs) — only idiosyncratic co-movement survives", font_size=15)
    _add_para(tf, "  - Factor groups by risk exposure — finds cross-sector pairs like NKE/YUM (consumer beta)", font_size=15)
    _add_para(tf, "  - Only 12% of pairs overlap across all three — each metric sees different structure", font_size=15)


def _pair_overlap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Pair Selection Overlap Across Metrics", font_size=36, bold=True, color=DARK_BLUE)

    img_path = ASSETS / "pair_overlap.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), Inches(9.4))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(1.8),
                      "Each metric captures different market structure:", font_size=18, bold=True)
    _add_para(tf, "  - SSD: 2,414 pairs unique to it — price-path similarity finds many candidates", font_size=15)
    _add_para(tf, "  - PC: 604 unique pairs — correlation-based filtering is the strictest", font_size=15)
    _add_para(tf, "  - Factor: 903 unique pairs — shared risk exposures find cross-sector relationships", font_size=15)
    _add_para(tf, "  - Only 307 pairs are found by all three metrics", font_size=15)


def _sector_composition_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "How Each Metric Clusters — Sector Composition", font_size=36, bold=True, color=DARK_BLUE)

    img_path = ASSETS / "sector_composition.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.2), Inches(1.0), Inches(9.6))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(0.5), Inches(5.6), Inches(9), Inches(2),
                      "Key difference:", font_size=18, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - PC is 91% same-sector — it finds idiosyncratic co-movement within industries", font_size=15)
    _add_para(tf, "  - SSD is only 57% same-sector — price-path similarity crosses sector boundaries", font_size=15)
    _add_para(tf, "  - Factor is 82% same-sector — shared risk factors mostly align with industry", font_size=15)
    _add_para(tf, "  - PC's selectivity (fewer, tighter pairs) explains its lower volatility and smoother returns", font_size=15)


def _top_trades_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Biggest Winners & Losers by Metric (IS)", font_size=36, bold=True, color=DARK_BLUE)

    img_path = ASSETS / "top_trades.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.1), Inches(1.0), Inches(9.8))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)

    tf = _add_textbox(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(1.8),
                      "Pattern:", font_size=18, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - SSD has the most extreme trades (+64% to -63%) — GFC bank pairs dominate both tails", font_size=14)
    _add_para(tf, "  - PC is more controlled (+43% to -40%) — GM/F (pre-bankruptcy) is the worst loss", font_size=14)
    _add_para(tf, "  - Factor is the most bounded (+28% to -27%) — fewer tail blowups", font_size=14)
    _add_para(tf, "  - Winners mostly revert; losers are almost all force-closes — confirming the bimodal thesis", font_size=14)


def _top_trades_oos_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Biggest Winners & Losers by Metric (OOS)", font_size=36, bold=True, color=DARK_BLUE)

    rows = [
        ["", "Winner", "Return", "Loser", "Return"],
        ["SSD", "BIIB/CNC '21", "+24.0%", "FRC/EL '23", "-61.2%"],
        ["", "CHTR/META '22", "+20.6%", "INTC/WDC '24", "-18.6%"],
        ["", "SBUX/VRSN '24", "+17.2%", "VFC/WYNN '22", "-18.3%"],
        ["PC", "EW/ZBH '24", "+21.5%", "GPC/LKQ '25", "-13.2%"],
        ["", "DOW/CE '24", "+12.6%", "MU/WDC '25", "-12.4%"],
        ["", "DOW/CE '25", "+11.9%", "FISV/FIS '23", "-12.3%"],
        ["Factor", "BSX/CNC '21", "+9.7%", "PTC/IT '25", "-19.4%"],
        ["", "WDC/QCOM '25", "+9.6%", "INTC/QCOM '21", "-15.3%"],
        ["", "VFC/MHK '23", "+9.2%", "NVDA/MPWR '24", "-14.9%"],
    ]
    _add_table(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(4.5), rows,
               col_widths=[Inches(1.2), Inches(2.4), Inches(1.2), Inches(2.4), Inches(1.2)])

    tf = _add_textbox(slide, Inches(0.5), Inches(5.9), Inches(9), Inches(1.6),
                      "OOS findings:", font_size=18, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - SSD's biggest OOS loss: FRC/EL (-61%) — First Republic Bank collapse, Mar 2023", font_size=14)
    _add_para(tf, "  - PC has the tightest loss range (-13%) — idiosyncratic selection avoids blow-ups", font_size=14)
    _add_para(tf, "  - Factor's biggest loss: NVDA/MPWR — AI momentum broke the factor-beta relationship", font_size=14)


def _carryover_intro_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Carry-Over: Letting Trades Roll Across Months", font_size=36, bold=True, color=DARK_BLUE)

    tf = _add_textbox(slide, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5),
                      "Instead of force-closing every position at month-end, we allow trades to roll over.",
                      font_size=20)
    _add_para(tf, "")
    _add_para(tf, "How it works:", font_size=20, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - Open position rolls into the next month if the pair is still in the candidate set", font_size=18)
    _add_para(tf, "  - Max hold: 3 months (aligned with half-life upper bound of 60 trading days)", font_size=18)
    _add_para(tf, "  - Position still exits on reversion (z=0) or at the 3-month cap", font_size=18)
    _add_para(tf, "")
    _add_para(tf, "Motivation:", font_size=20, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - ~89% of trades in the baseline are force-closed at month-end", font_size=18)
    _add_para(tf, "  - Most of those positions haven't reverted yet — they just need more time", font_size=18)
    _add_para(tf, "  - Force-close drag is the #1 source of lost profit", font_size=18)
    _add_para(tf, "")
    _add_para(tf, "All results below use cointegration-filtered strategies with carry_over=True, max_carry_months=3.",
              font_size=16, color=GREY)


def _carryover_results_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Carry-Over Results — IS vs OOS Scorecard", font_size=36, bold=True, color=DARK_BLUE)

    rows = [
        ["Strategy", "Period", "Ann. Ret", "Ann. Vol", "Sharpe", "Max DD", "Months"],
        ["SSD + filter", "IS 2003-2020", "", "", "0.889", "", "215"],
        ["Factor + filter", "IS 2003-2020", "", "", "0.862", "", "215"],
        ["PC + filter", "IS 2003-2020", "", "", "0.736", "", "215"],
        ["SSD + filter", "OOS 2021-2025", "", "", "0.542", "", "59"],
        ["PC + filter", "OOS 2021-2025", "", "", "0.183", "", "59"],
        ["Factor + filter", "OOS 2021-2025", "", "", "-0.056", "", "59"],
    ]

    # Fill in actual metrics
    import pandas as pd
    import sys
    ROOT = Path(__file__).resolve().parent.parent.parent
    sys.path.insert(0, str(ROOT))
    try:
        from src.performance import compute_metrics
        fmap = {
            1: "carry_ssd_filtered", 2: "carry_factor_filtered", 3: "carry_pc_filtered",
            4: "carry_oos_ssd_filtered", 5: "carry_oos_pc_filtered", 6: "carry_oos_factor_filtered",
        }
        for row_idx, fname in fmap.items():
            path = Path(__file__).resolve().parent.parent / "results" / f"{fname}_monthly.parquet"
            if path.exists():
                df = pd.read_parquet(path)
                m = compute_metrics(df["monthly_return"])
                rows[row_idx][2] = f"{m.ann_return:.2%}"
                rows[row_idx][3] = f"{m.ann_vol:.2%}"
                rows[row_idx][4] = f"{m.sharpe:.3f}"
                rows[row_idx][5] = f"{m.max_drawdown:.2%}"
                rows[row_idx][6] = str(m.n_months)
    except Exception:
        pass

    _add_table(slide, Inches(0.3), Inches(1.2), Inches(9.4), Inches(3.8), rows)

    tf = _add_textbox(slide, Inches(0.5), Inches(5.3), Inches(9), Inches(2.2),
                      "Highlights:", font_size=18, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - SSD + filter is the strongest carry-over strategy: 0.889 IS, 0.542 OOS", font_size=16)
    _add_para(tf, "  - Factor + filter IS stays strong (0.862) but turns negative OOS (-0.056)", font_size=16)
    _add_para(tf, "  - PC + filter weakens in both IS (0.736) and OOS (0.183)", font_size=16)


def _carryover_cumulative_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Cumulative Returns — Carry-Over Strategies (IS + OOS)", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "cumulative_pnl_carryover.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.3), Inches(1.0), Inches(9.4))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(9), Inches(0.6),
                 "Solid lines = IS (2003-2020). Dashed lines = OOS (2021-2025). All with carry_over=True, max 3 months.",
                 font_size=14, color=GREY)


def _carryover_yearly_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Year-by-Year Returns — Carry-Over Strategies", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "yearly_returns_carryover.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.2), Inches(1.0), Inches(9.6))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    tf = _add_textbox(slide, Inches(0.5), Inches(5.5), Inches(9), Inches(2),
                      "Findings:", font_size=16, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - SSD + filter carry-over: consistently positive across years, strong in GFC and OOS", font_size=14)
    _add_para(tf, "  - More time for reversion smooths out year-to-year variance for SSD", font_size=14)
    _add_para(tf, "  - PC and factor show more negative years when trades are held longer", font_size=14)


def _carryover_pnl_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "P&L Breakdown — Carry-Over Strategies (IS vs OOS)", font_size=34, bold=True, color=DARK_BLUE)
    img_path = ASSETS / "pnl_breakdown_carryover.png"
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.1), Inches(1.0), Inches(9.8))
    else:
        _add_textbox(slide, Inches(2), Inches(3), Inches(6), Inches(1),
                     "[Chart will appear after running: python submission/report/build_charts.py]",
                     font_size=16, color=GREY, align=PP_ALIGN.CENTER)
    tf = _add_textbox(slide, Inches(0.5), Inches(5.8), Inches(9), Inches(1.8),
                      "Key change vs force-close:", font_size=16, bold=True, color=DARK_BLUE)
    _add_para(tf, "  - Higher reversion rate — more trades complete their mean-reversion cycle", font_size=14)
    _add_para(tf, "  - Fewer force-closes — the drag that was killing profitability is reduced", font_size=14)
    _add_para(tf, "  - Trade-off: trades held longer carry more risk per position", font_size=14)


def _carryover_findings_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 "Carry-Over Key Findings", font_size=36, bold=True, color=DARK_BLUE)

    tf = _add_textbox(slide, Inches(0.5), Inches(1.1), Inches(9), Inches(6),
                      "", font_size=18)
    _add_para(tf, "1. SSD + filter benefits most from carry-over", font_size=22, bold=True, color=DARK_BLUE)
    _add_para(tf, "     IS Sharpe 0.889  |  OOS Sharpe 0.542", font_size=18)
    _add_para(tf, "     SSD selects by price-path shape — when these pairs diverge, they genuinely", font_size=16)
    _add_para(tf, "     need more time to converge. Giving them 3 months lets the reversion complete.", font_size=16)
    _add_para(tf, "")
    _add_para(tf, "2. PC + filter works better with force-close", font_size=22, bold=True, color=DARK_BLUE)
    _add_para(tf, "     IS Sharpe 0.736  |  OOS Sharpe 0.183", font_size=18)
    _add_para(tf, "     PC pairs are selected by idiosyncratic correlation — when they haven't reverted", font_size=16)
    _add_para(tf, "     by month-end, the correlation structure may have shifted. Cutting losses early is better.", font_size=16)
    _add_para(tf, "")
    _add_para(tf, "3. Factor-beta also prefers force-close", font_size=22, bold=True, color=DARK_BLUE)
    _add_para(tf, "     IS Sharpe 0.862  |  OOS Sharpe -0.056", font_size=18)
    _add_para(tf, "     Factor loadings shift over time — holding amplifies the mismatch.", font_size=16)
    _add_para(tf, "")
    _add_para(tf, "Takeaway: the optimal trade duration depends on the distance metric.",
              font_size=20, bold=True, color=DARK_BLUE)
    _add_para(tf, "SSD pairs revert slowly (price convergence); PC/factor pairs revert quickly or not at all.",
              font_size=18)


def _closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)
    _add_textbox(slide, Inches(1), Inches(2), Inches(8), Inches(1),
                 "Thank You", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.6),
                 "Clustering-Based Pairs Trading", font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.6),
                 "IS Sharpe 1.113  |  OOS Sharpe 0.412  |  274 months  |  built from the data up",
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.5),
                 "Happy to walk through the code, the backtest engine, or the attribution analysis.",
                 font_size=14, color=GREY, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)
    _problem_slide(prs)
    _data_slide(prs)
    _pipeline_slide(prs)
    _distance_slide(prs)
    _clustering_slide(prs)
    _signal_slide(prs)
    _factor_beta_slide(prs)
    _is_results_slide(prs)
    _cumulative_pnl_slide(prs)
    _attribution_slide(prs)
    _pnl_breakdown_slide(prs)
    _yearly_returns_slide(prs)
    _oos_slide(prs)
    _filtered_results_slide(prs)
    _cumulative_pnl_filtered_slide(prs)
    _yearly_returns_filtered_slide(prs)
    _pnl_breakdown_filtered_slide(prs)
    _pair_examples_slide(prs)
    _pair_overlap_slide(prs)
    _sector_composition_slide(prs)
    _top_trades_slide(prs)
    _top_trades_oos_slide(prs)
    _carryover_intro_slide(prs)
    _carryover_results_slide(prs)
    _carryover_cumulative_slide(prs)
    _carryover_yearly_slide(prs)
    _carryover_pnl_slide(prs)
    _carryover_findings_slide(prs)
    _rigour_slide(prs)
    _takeaway_slide(prs)
    _closing_slide(prs)

    prs.save(str(OUT))
    print(f"Saved {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
